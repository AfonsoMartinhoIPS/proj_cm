from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_BG      = RGBColor(0x0F, 0x17, 0x2A)
C_ACCENT  = RGBColor(0x00, 0xC8, 0x8E)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xCC, 0xE8, 0xFF)
C_SUB     = RGBColor(0x99, 0xBB, 0xDD)
C_DARK    = RGBColor(0x1A, 0x28, 0x45)
C_ROW1    = RGBColor(0x14, 0x22, 0x3A)
C_ROW2    = RGBColor(0x1A, 0x2B, 0x47)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = C_BG

def tb(slide, text, x, y, w, h, size=18, bold=False, color=C_WHITE,
       align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tf

def accent_bar(slide, y=Inches(0.55), w=Inches(1.2)):
    s = slide.shapes.add_shape(1, Inches(0.5), y, w, Pt(4))
    s.fill.solid(); s.fill.fore_color.rgb = C_ACCENT; s.line.fill.background()

def slide_header(slide, section, title, subtitle=None):
    accent_bar(slide)
    tb(slide, section, Inches(0.5), Inches(0.35), Inches(5), Inches(0.38),
       size=11, bold=True, color=C_ACCENT)
    tb(slide, title, Inches(0.5), Inches(0.65), Inches(12.3), Inches(0.72),
       size=32, bold=True)
    if subtitle:
        tb(slide, subtitle, Inches(0.5), Inches(1.38), Inches(12.3), Inches(0.45),
           size=15, color=C_SUB)

def bullets_box(slide, items, x, y, w, h, size=15, color=C_LIGHT, bullet="•"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color

def card(slide, title, lines, x, y, w=Inches(3.8), h=Inches(2.2), title_size=14, body_size=12):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = C_DARK
    sh.line.color.rgb = C_ACCENT; sh.line.width = Pt(1)
    tb(slide, title, x+Inches(0.15), y+Inches(0.1), w-Inches(0.3), Inches(0.42),
       size=title_size, bold=True, color=C_ACCENT)
    bullets_box(slide, lines, x+Inches(0.15), y+Inches(0.55),
                w-Inches(0.3), h-Inches(0.65), size=body_size)

def pill(slide, text, x, y, w=Inches(2.8), h=Inches(0.5)):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x00, 0x50, 0x3A)
    sh.line.color.rgb = C_ACCENT; sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = text; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = C_ACCENT

def dot_point(slide, title, desc, x, y, w=Inches(5.5)):
    dot = slide.shapes.add_shape(1, x, y+Inches(0.08), Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = C_ACCENT; dot.line.fill.background()
    tb(slide, title, x+Inches(0.35), y, w, Inches(0.38), size=14, bold=True, color=C_ACCENT)
    tb(slide, desc,  x+Inches(0.35), y+Inches(0.38), w, Inches(0.75), size=13, color=C_LIGHT)

# ─── SLIDE 1: CAPA ──────────────────────────────────────────────────────────
s = add_slide(); bg(s)
bar = s.shapes.add_shape(1, Inches(0.5), Inches(2.1), Inches(0.08), Inches(2.4))
bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT; bar.line.fill.background()

tb(s, "NutriScan", Inches(0.8), Inches(2.0), Inches(11), Inches(1.3),
   size=68, bold=True)
tb(s, "Rastreio nutricional e comparação de preços em supermercado",
   Inches(0.8), Inches(3.3), Inches(10), Inches(0.65), size=20, color=C_SUB)

div = s.shapes.add_shape(1, Inches(0.8), Inches(4.15), Inches(5), Pt(1))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0x1E,0x35,0x55); div.line.fill.background()

tb(s, "Computação Móvel  |  2025/26", Inches(0.8), Inches(4.3), Inches(8), Inches(0.4),
   size=14, color=RGBColor(0x55,0x77,0x99))
tb(s, "Prof. Bruno Pereira  |  Prof. Paula Miranda", Inches(0.8), Inches(4.75),
   Inches(8), Inches(0.4), size=13, color=RGBColor(0x44,0x66,0x88))

students = "Samuel Silva (202200315)   •   Afonso Martinho (202001865)\nDaniel Pais (202200286)   •   Fernando Ramalho (202002203)"
tb(s, students, Inches(0.8), Inches(5.4), Inches(11), Inches(0.8), size=13, color=C_SUB)

# ─── SLIDE 2: CONCEITO ──────────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_header(s, "01  CONCEITO", "O que é o NutriScan?")

tb(s, "Aplicação móvel Flutter que une dois domínios habitualmente separados:",
   Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.5), size=17, color=C_LIGHT)

card(s, "Rastreio Nutricional",
     ["Digitalizar produtos pelo código de barras",
      "Ficha nutricional completa",
      "Registar refeições e água diariamente",
      "Acompanhar progresso face a objetivos"],
     Inches(0.5), Inches(2.55), w=Inches(5.9), h=Inches(2.3))

card(s, "Comparação de Preços",
     ["Registar preço visto em qualquer loja",
      "Cálculo automático de preço por kg/L",
      "Comparar opções entre lojas e marcas",
      "Destaque da opção mais barata"],
     Inches(6.9), Inches(2.55), w=Inches(5.9), h=Inches(2.3))

pb = s.shapes.add_shape(1, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.75))
pb.fill.solid(); pb.fill.fore_color.rgb = RGBColor(0x00,0x30,0x22)
pb.line.color.rgb = C_ACCENT; pb.line.width = Pt(1)
tb(s, "Problema: apps de nutrição não têm preços  •  Usar inúmeras apps para atingir um objetivo único  •  NutriScan resolve ambos",
   Inches(0.65), Inches(5.2), Inches(12.0), Inches(0.55),
   size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ─── SLIDE 3: PÚBLICO-ALVO ──────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_header(s, "02  PÚBLICO-ALVO", "Quem usa o NutriScan?",
             "Pessoas que querem comer melhor e gastar menos no supermercado")

profiles = [
    ("Jovens e Estudantes",     ["Orçamento controlado", "Quer escolhas mais inteligentes", "Habituados a apps móveis"]),
    ("Objetivos de Saúde",      ["Perda de peso", "Ganho de massa muscular", "Manutenção de saúde"]),
    ("Compradores Conscientes", ["Fazem as próprias compras", "Querem decisões mais informadas", "Comparam preços manualmente hoje"]),
    ("Dietas Específicas",      ["Controlo rigoroso de macros", "Precisam de saber exatamente o que comem", "Contagem calórica diária"]),
]
for i, (title, lines) in enumerate(profiles):
    x = Inches(0.5 + i * 3.2)
    card(s, title, lines, x, Inches(2.3), w=Inches(2.9), h=Inches(2.3))

# ─── SLIDE 4: FUNCIONALIDADES GERAIS ────────────────────────────────────────
s = add_slide(); bg(s)
slide_header(s, "03  FUNCIONALIDADES GERAIS", "Funcionalidades comuns a qualquer app móvel")

left = [
    "Splash Screen – verificação de sessão no arranque",
    "Onboarding – wizard antes do registo (dados, objetivo, cálculo TMB)",
    "Autenticação – registo e login via Firebase Auth",
    "Sessão persistente – sem re-login entre sessões",
]
right = [
    "Navegação – bottom nav 4 tabs + stack para detalhes",
    "Estado das tabs preservado ao alternar",
    "Perfil e Definições – objetivos, notificações, logout",
    "Créditos – autores, UC, ano letivo",
]
bullets_box(s, left,  Inches(0.5), Inches(2.1), Inches(6.1), Inches(3.5), size=15)
bullets_box(s, right, Inches(6.9), Inches(2.1), Inches(6.1), Inches(3.5), size=15)

div = s.shapes.add_shape(1, Inches(6.6), Inches(2.0), Pt(2), Inches(3.8))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0x1E,0x35,0x55); div.line.fill.background()

# ─── SLIDE 5: ONBOARDING ────────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_header(s, "03  FUNCIONALIDADES GERAIS", "Onboarding",
             "Primeira interação — acontece antes do registo, \"agarra\" logo o utilizador")

steps = [
    ("1. Dados Pessoais",   "Nome, idade, peso,\naltura, sexo"),
    ("2. Objetivo",         "Perder peso /\nManter / Ganhar massa"),
    ("3. Cálculo Auto",     "Fórmula Mifflin-St Jeor\nTMB + fator atividade"),
    ("4. Confirmar",        "Utilizador revê e\najusta os valores"),
]
for i, (step, desc) in enumerate(steps):
    x = Inches(0.5 + i * 3.2)
    pill(s, step, x, Inches(2.2))
    tb(s, desc, x, Inches(2.85), Inches(2.8), Inches(1.1), size=14, color=C_LIGHT)
    if i < 3:
        tb(s, ">>", Inches(3.1 + i * 3.2), Inches(2.37), Inches(0.4), Inches(0.38),
           size=16, bold=True, color=C_ACCENT)

tb(s, "Fórmula TMB (Mifflin-St Jeor)", Inches(0.5), Inches(4.2), Inches(12), Inches(0.4),
   size=14, bold=True, color=C_ACCENT)
tb(s, "Homem:  TMB = 10 × peso + 6.25 × altura − 5 × idade + 5\n"
      "Mulher: TMB = 10 × peso + 6.25 × altura − 5 × idade − 161\n"
      "Calorias = TMB × fator de atividade  (±300 kcal por objetivo)",
   Inches(0.5), Inches(4.65), Inches(12), Inches(1.2), size=13, color=C_LIGHT)

tb(s, "Macros:  Proteína = 2.0g × peso  |  Gordura = 25% calorias ÷ 9  |  Hidratos = restante ÷ 4  |  Água = 35ml × peso",
   Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.4),
   size=12, color=C_SUB, align=PP_ALIGN.CENTER)

# ─── SLIDE 6: PRODUTOS ──────────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_header(s, "04  FUNCIONALIDADES ESPECÍFICAS", "Produtos",
             "Repositório de alimentos do utilizador")

card(s, "Scanner (Código de Barras)",
     ["Câmara em tempo real",
      "Consulta OpenFoodFacts por barcode",
      "Nome, marca, possivelmente imagem",
      "Tabela nutricional por 100g/ml",
      "Guardar produto para uso futuro"],
     Inches(0.5), Inches(2.1), w=Inches(3.8), h=Inches(2.6))

card(s, "Pesquisa por Nome",
     ["Para alimentos sem código de barras",
      "Fonte 1: OpenFoodFacts (embalados)",
      "Fonte 2: USDA FoodData Central",
      "Fallback automático OFF → USDA",
      "Resultado no mesmo modelo de dados"],
     Inches(4.7), Inches(2.1), w=Inches(3.8), h=Inches(2.6))

card(s, "Registo de Preços em Loja",
     ["Registar preço (€) + loja + quantidade",
      "Cálculo automático preço/kg ou /L",
      "Comparação entre lojas e embalagens",
      "Destaque da opção mais económica"],
     Inches(8.9), Inches(2.1), w=Inches(3.9), h=Inches(2.6))

card(s, "Lista de Produtos Guardados",
     ["Todos os produtos já digitalizados ou pesquisados",
      "Pesquisa local para filtrar rapidamente",
      "Indicador de fonte: OFF / USDA / Manual"],
     Inches(0.5), Inches(4.95), w=Inches(12.3), h=Inches(1.5))

# ─── SLIDE 7: REFEIÇÕES E HISTÓRICO ────────────────────────────────────────
s = add_slide(); bg(s)
slide_header(s, "04  FUNCIONALIDADES ESPECÍFICAS", "Refeições e Histórico",
             "Diário alimentar e visão retrospetiva")

card(s, "Registo de Refeições",
     ["PA, Almoço, Jantar, Snack",
      "Quantidade consumida (g ou ml)",
      "Cálculo nutricional por porção",
      "Barras de progresso face a objetivos"],
     Inches(0.5), Inches(2.1), w=Inches(5.9), h=Inches(2.2))

card(s, "Registo de Água",
     ["Botões rápidos: 200ml, 330ml, 500ml",
      "Campo manual para valor personalizado",
      "Barra de progresso face ao objetivo"],
     Inches(0.5), Inches(4.5), w=Inches(5.9), h=Inches(1.8))

card(s, "Histórico – Vista Dia",
     ["Navegar entre dias ou por calendário",
      "Ver refeições + totais nutricionais",
      "Comparação com objetivos"],
     Inches(6.9), Inches(2.1), w=Inches(5.9), h=Inches(1.7))

card(s, "Histórico – Vista Semana",
     ["Gráfico de barras: calorias por dia (7d)",
      "Médias semanais por macro",
      "Destaque dias acima/abaixo do objetivo"],
     Inches(6.9), Inches(4.0), w=Inches(5.9), h=Inches(1.7))

card(s, "Histórico – Vista Mês",
     ["Calendário com cor por dia (verde/vermelho/cinza)",
      "Resumo mensal: média diária, dias com registo"],
     Inches(6.9), Inches(5.9), w=Inches(5.9), h=Inches(1.4))

# ─── SLIDE 8: ORIGINALIDADE ─────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_header(s, "05  ORIGINALIDADE E MAIS-VALIAS", "O que nos diferencia?")

points = [
    ("Nutrição + preços na mesma app",
     "Um local só, organizado e fácil de utilizar"),
    ("Dupla fonte de dados nutricional",
     "OpenFoodFacts (embalados) + USDA (genéricos)"),
    ("Onboarding com cálculo científico",
     "Objetivos calculados com fórmula Mifflin-St Jeor, não valores genéricos"),
    ("Registo de água integrado",
     "Visão completa de ingestão diária, não apenas alimentos sólidos"),
    ("Preço por kg/L automático",
     "Comparação justa entre embalagens de tamanhos diferentes"),
    ("Histórico multi-escala",
     "Dia / semana / mês: identifica padrões ao longo do tempo"),
]
for i, (title, desc) in enumerate(points):
    row = i // 2
    col = i % 2
    dot_point(s, title, desc, Inches(0.5 + col * 6.4), Inches(2.0 + row * 1.55))

# ─── SLIDE 9: NAVEGAÇÃO ──────────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_header(s, "06  ESQUEMA DE NAVEGAÇÃO", "Fluxo entre ecrãs")

nav = (
    "SplashScreen\n"
    "  |-- (1ª vez)        -->  Onboarding (4 steps)  -->  Register  -->  Login  -->  Home\n"
    "  |-- (sessão ativa)  -->  Home\n"
    "\n"
    "Home  -->  Perfil  -->  Definições  -->  Créditos\n"
    "Home  -->  Histórico  (Dia / Semana / Mês)\n"
    "\n"
    "Refeições  -->  Adicionar  -->  Pesquisa / Scanner / Guardados  -->  ProductDetail\n"
    "\n"
    "Produtos   -->  ProductDetail  -->  Registar Preço  -->  Lista de Preços\n"
    "\n"
    "Scan (tab) -->  ScannerScreen  -->  ProductDetail\n"
    "                |-- Inserir manualmente  -->  SearchScreen  -->  ProductDetail"
)
box = s.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(5.0))
tf = box.text_frame; tf.word_wrap = False
for i, line in enumerate(nav.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = line
    r.font.name = "Courier New"
    r.font.size = Pt(14)
    top = line.strip().startswith(("SplashScreen","Home","Refeições","Produtos","Scan"))
    r.font.color.rgb = C_ACCENT if top else C_LIGHT
    r.font.bold = top

# ─── SLIDE 10: STACK TÉCNICA ─────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_header(s, "07  STACK TÉCNICA", "Tecnologias utilizadas")

tech = [
    ("Framework",          "Flutter (Dart)"),
    ("Gestão de estado",   "Riverpod (ex.)"),
    ("Autenticação",       "Firebase Authentication"),
    ("Base de dados",      "Cloud Firestore"),
    ("API embalados",      "OpenFoodFacts"),
    ("API genéricos",      "USDA FoodData Central"),
    ("Barcode scanner",    "mobile_scanner"),
    ("HTTP",               "dio (ex.)"),
    ("Notificações",       "flutter_local_notifications"),
    ("Gráficos",           "fl_chart (ex.)"),
    ("Arquitetura",        "Clean Arch. – data / domain / presentation"),
]
col_n = 6
for i, (comp, tech_name) in enumerate(tech):
    row = i % col_n
    col = i // col_n
    x = Inches(0.5 + col * 6.4)
    y = Inches(2.05 + row * 0.78)
    rect = s.shapes.add_shape(1, x, y, Inches(5.9), Inches(0.62))
    rect.fill.solid()
    rect.fill.fore_color.rgb = C_ROW1 if i % 2 == 0 else C_ROW2
    rect.line.fill.background()
    tb(s, comp,      x+Inches(0.15), y+Inches(0.1), Inches(2.0), Inches(0.42), size=13, color=C_SUB)
    tb(s, tech_name, x+Inches(2.2),  y+Inches(0.1), Inches(3.5), Inches(0.42), size=13, bold=True)

out = "/Users/samuelsilva/Library/Mobile Documents/com~apple~CloudDocs/personal/Universidade/2025_26/Semestre2/CM/projeto/NutriScan_Apresentacao.pptx"
prs.save(out)
print("Saved:", out)
